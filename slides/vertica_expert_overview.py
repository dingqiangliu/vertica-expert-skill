#!/usr/bin/env python3
"""
Vertica Expert Skill - Comprehensive Overview Presentation
Generates a professional PowerPoint presentation introducing the Vertica Expert skill.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
import os

# ── Color Palette ──────────────────────────────────────────────────────────
DARK_BLUE = RGBColor(0x1B, 0x2A, 0x4A)
MID_BLUE = RGBColor(0x2E, 0x50, 0x80)
ACCENT_BLUE = RGBColor(0x4A, 0x90, 0xD9)
ORANGE = RGBColor(0xF5, 0x8C, 0x00)
GREEN = RGBColor(0x28, 0xA7, 0x45)
LIGHT_GRAY = RGBColor(0xF0, 0xF2, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x1B, 0x2A, 0x4A)
MEDIUM_TEXT = RGBColor(0x4A, 0x55, 0x68)
TEAL = RGBColor(0x17, 0xA2, 0xB8)
PURPLE = RGBColor(0x6F, 0x42, 0xC1)
RED = RGBColor(0xE7, 0x4C, 0x3C)
SOFT_BLUE = RGBColor(0xD6, 0xEA, 0xF8)
SOFT_GREEN = RGBColor(0xD4, 0xED, 0xDA)
SOFT_ORANGE = RGBColor(0xFF, 0xF3, 0xCD)
SOFT_PURPLE = RGBColor(0xE8, 0xDA, 0xEF)
SOFT_RED = RGBColor(0xF8, 0xD7, 0xDA)
SOFT_TEAL = RGBColor(0xD1, 0xEC, 0xF1)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


# ── Helper Functions ────────────────────────────────────────────────────────

def add_background(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, border_color=None, border_width=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        if border_width:
            shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=14, color=DARK_TEXT, bold=False,
                 alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=13, color=DARK_TEXT,
                    spacing=Pt(6), font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = spacing
        p.level = 0
    return txBox


def add_code_block(slide, left, top, width, height, code_text, font_size=8.5,
                   bg_color=RGBColor(0x1E, 0x1E, 0x2E), text_color=RGBColor(0xD4, 0xD4, 0xD4)):
    """Add a code block with dark background."""
    bg = add_shape(slide, left, top, width, height, bg_color, border_color=RGBColor(0x33, 0x33, 0x33), border_width=Pt(1))
    txBox = slide.shapes.add_textbox(left + Inches(0.1), top + Inches(0.05), width - Inches(0.2), height - Inches(0.1))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = code_text
    p.font.size = Pt(font_size)
    p.font.color.rgb = text_color
    p.font.name = "Consolas"
    return txBox


def add_title_bar(slide, title_text, subtitle_text=None):
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.0), DARK_BLUE)
    add_shape(slide, Inches(0), Inches(1.0), Inches(13.333), Inches(0.06), ORANGE)
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.08), Inches(12.3), Inches(0.6))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(26)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = "Segoe UI"
    if subtitle_text:
        p2 = tf.add_paragraph()
        p2.text = subtitle_text
        p2.font.size = Pt(13)
        p2.font.color.rgb = ACCENT_BLUE
        p2.font.name = "Segoe UI"


def add_footer(slide, text="Vertica Expert Skill  |  Comprehensive Migration & Development Guide"):
    add_shape(slide, Inches(0), Inches(7.15), Inches(13.333), Inches(0.35), MID_BLUE)
    txBox = slide.shapes.add_textbox(Inches(0.3), Inches(7.18), Inches(12.7), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(0xAA, 0xBB, 0xCC)
    p.font.name = "Segoe UI"
    p.alignment = PP_ALIGN.RIGHT


def add_icon_circle(slide, left, top, size_w, size_h, color, label=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size_w, size_h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if label is not None:
        txBox = slide.shapes.add_textbox(left, top + Inches(0.05), size_w, size_h)
        tf = txBox.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = str(label)
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.font.name = "Segoe UI"
        p.alignment = PP_ALIGN.CENTER
    return shape


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title Slide
# ═══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, DARK_BLUE)

add_shape(slide, Inches(0), Inches(0), Inches(4), Inches(7.5), MID_BLUE)
add_shape(slide, Inches(0), Inches(6.5), Inches(13.333), Inches(1.0), RGBColor(0x15, 0x22, 0x3A))

for pos in [(10.5, 1.0), (11.5, 2.5), (11.0, 4.0), (12.2, 1.8)]:
    c = add_shape(slide, Inches(pos[0]), Inches(pos[1]), Inches(0.8), Inches(0.8), RGBColor(0x4A, 0x90, 0xD9))
    c.fill.background()
    c.line.color.rgb = ACCENT_BLUE
    c.line.width = Pt(2)

add_text_box(slide, Inches(1.2), Inches(1.5), Inches(8), Inches(1.5), "Vertica Expert",
             font_size=54, color=WHITE, bold=True)
add_text_box(slide, Inches(1.2), Inches(3.0), Inches(8), Inches(1.0),
             "Comprehensive Skill for Database Migration & Development",
             font_size=24, color=ACCENT_BLUE)
add_text_box(slide, Inches(1.2), Inches(4.3), Inches(8), Inches(1.5),
             "SQL Syntax  •  Stored Procedures  •  UDx  •  Machine Learning\nOracle / DB2 / SQL Server / PostgreSQL / MySQL → Vertica",
             font_size=16, color=RGBColor(0xAA, 0xBB, 0xCC))

badge_x = Inches(9.5)
badges = [
    ("5", "Migration Paths", ORANGE),
    ("100+", "Function Mappings", GREEN),
    ("7", "ML Algorithms", TEAL),
    ("18", "Reference Guides", PURPLE),
]
for i, (num, label, color) in enumerate(badges):
    y_pos = Inches(1.2 + i * 1.3)
    add_shape(slide, badge_x - Inches(0.7), y_pos + Inches(0.05), Inches(0.6), Inches(0.6), color)
    add_text_box(slide, badge_x - Inches(0.7), y_pos + Inches(0.05), Inches(0.6), Inches(0.6), num,
                font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, badge_x, y_pos, Inches(3.5), Inches(0.8), label, font_size=16, color=WHITE)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — What is Vertica Expert?
# ═══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_title_bar(slide, "What is the Vertica Expert Skill?", "A comprehensive toolkit for Vertica migration & development")

add_text_box(slide, Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.8),
             "The Vertica Expert skill is an AI-powered knowledge base that provides comprehensive guidance for migrating databases from other systems to Vertica and optimizing SQL workflows for Vertica's columnar MPP architecture.",
             font_size=13, color=DARK_TEXT)

capabilities = [
    ("🔄", "Database Migration", "Convert DDL, DML, stored procedures, and queries from Oracle, DB2, SQL Server, PostgreSQL, and MySQL to native Vertica syntax"),
    ("⚡", "Performance Optimization", "Design projections, encoding strategies, and resource management for maximum columnar performance"),
    ("🤖", "Machine Learning", "Implement end-to-end in-database ML workflows — regression, classification, clustering, and time series"),
    ("🔧", "UDx Development", "Build custom scalar, aggregate, analytic, and transform functions in C++, Python, Java, or R"),
    ("📊", "SQL Development", "Write complex analytical queries using Vertica's advanced SQL syntax, window functions, and CTEs"),
    ("📋", "VSQL Testing", "Immediately test all converted SQL and stored procedures using the VSQL command-line framework"),
]

for i, (icon, title, desc) in enumerate(capabilities):
    col = i % 2
    row = i // 2
    x = Inches(0.5 + col * 6.3)
    y = Inches(2.2 + row * 1.55)
    card = add_shape(slide, x, y, Inches(6.1), Inches(1.4), SOFT_BLUE, border_color=ACCENT_BLUE, border_width=Pt(1))
    add_text_box(slide, x + Inches(0.15), y + Inches(0.08), Inches(5.8), Inches(0.4),
                 f"{icon}  {title}", font_size=14, color=DARK_BLUE, bold=True)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.48), Inches(5.8), Inches(0.85),
                 desc, font_size=10.5, color=MEDIUM_TEXT)

add_footer(slide)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Migration Paths Supported
# ═══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_title_bar(slide, "Supported Migration Paths", "Five major database systems → Vertica")

paths = [
    ("Oracle", "→ Vertica", "PL/SQL → PL/vSQL\nPackage migration\nSequence & trigger conversion", ORANGE, SOFT_ORANGE),
    ("IBM DB2", "→ Vertica", "PL/SQL → PL/vSQL\nModule/package conversion\nMQT → Live Aggregate Projections", GREEN, SOFT_GREEN),
    ("SQL Server", "→ Vertica", "T-SQL → Vertica SQL\nIdentity columns\nCursor alternatives", TEAL, SOFT_TEAL),
    ("PostgreSQL", "→ Vertica", "PL/pgSQL → PL/vSQL\nArray & JSON handling\nFunction mapping", PURPLE, SOFT_PURPLE),
    ("MySQL", "→ Vertica", "Schema & query conversion\nAUTO_INCREMENT → IDENTITY\nStorage engine mapping", RED, SOFT_RED),
]

for i, (source, arrow, details, accent, bg) in enumerate(paths):
    x = Inches(0.5 + i * 2.55)
    y = Inches(1.4)
    w = Inches(2.4)
    card = add_shape(slide, x, y, w, Inches(5.2), bg, border_color=accent, border_width=Pt(2))
    add_text_box(slide, x, y + Inches(0.3), w, Inches(0.5), source,
                font_size=18, color=DARK_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, y + Inches(0.75), w, Inches(0.4), arrow,
                font_size=20, color=accent, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.15), y + Inches(1.3), w - Inches(0.3), Inches(3.5),
                details, font_size=11, color=MEDIUM_TEXT, alignment=PP_ALIGN.CENTER)

add_footer(slide)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Migration Workflow
# ═══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_title_bar(slide, "Migration Workflow", "Step-by-step process for database migration")

steps = [
    ("1", "Read All References", "Read Generic Migration Guide (mandatory), OLTP→OLAP Rewrite Guide, and source-specific guide before touching any source file.", ORANGE),
    ("2", "List Requirements", "Present all understood requirements to the user and wait for confirmation before starting.", GREEN),
    ("3", "Sequential Processing", "Process source files in alphabetical order, one file at a time, top to bottom within each file.", TEAL),
    ("4", "Migrate & Rewrite", "Convert objects one-to-one (tables→tables, views→views, etc.) and rewrite OLTP patterns to set-based OLAP style.", PURPLE),
    ("5", "Test-First Rule", "MIGRATE → TEST → PASS → APPEND. Test every object immediately after migration. Never append untested code.", RED),
    ("6", "Validate & Report", "Run complete integration test, check error logs, and generate a comprehensive migration report.", DARK_BLUE),
]

for i, (num, title, desc, color) in enumerate(steps):
    y = Inches(1.3 + i * 0.95)
    add_icon_circle(slide, Inches(0.7), y + Inches(0.1), Inches(0.5), Inches(0.5), color, label=num)
    if i < len(steps) - 1:
        add_shape(slide, Inches(0.9), y + Inches(0.6), Inches(0.04), Inches(0.38), color)
    card = add_shape(slide, Inches(1.4), y, Inches(11.4), Inches(0.88), LIGHT_GRAY, border_color=color, border_width=Pt(1))
    add_text_box(slide, Inches(1.6), y + Inches(0.05), Inches(11), Inches(0.35),
                title, font_size=14, color=color, bold=True)
    add_text_box(slide, Inches(1.6), y + Inches(0.38), Inches(11), Inches(0.4),
                desc, font_size=10.5, color=MEDIUM_TEXT)

add_footer(slide)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Migration Rules & Prohibitions
# ═══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_title_bar(slide, "Mandatory Rules & Prohibitions", "Non-negotiable requirements for all migrations")

add_text_box(slide, Inches(0.5), Inches(1.3), Inches(5.8), Inches(0.4),
             "✅ Mandatory Rules", font_size=16, color=GREEN, bold=True)

rules = [
    "Migrate ALL objects — no selective migration",
    "Process files in alphabetical order — one at a time",
    "Process each file top-to-bottom — never skip or reorder",
    "Keep objects intact — never split procedures or statements",
    "One-to-one mapping — tables→tables, views→views, etc.",
    "Rewrite OLTP→OLAP — eliminate cursors, row-by-row DML",
    "Preserve ALL logic — never simplify or remove code",
    "Test EVERY object immediately — no exceptions",
    "Never use scripts/tools for bulk conversion",
    "Never modify original file ordering",
]

add_bullet_list(slide, Inches(0.5), Inches(1.75), Inches(5.8), Inches(4.8), rules, font_size=11.5, color=DARK_TEXT)

add_text_box(slide, Inches(6.8), Inches(1.3), Inches(6), Inches(0.4),
             "🚫 Absolutely Prohibited", font_size=16, color=RED, bold=True)

prohibitions = [
    ("Skipping objects", "\"Seem unnecessary\" objects must still be migrated"),
    ("Reordering", "Dependencies are already correctly ordered in source files"),
    ("Automated scripts", "Bulk processing loses business logic"),
    ("Batch processing", "Migrate and test one object at a time"),
    ("Removing OUT/INOUT", "Parameter keywords must be preserved"),
    ("Creating projections from indexes", "Comment out index statements instead"),
    ("Assuming unsupported", "Always verify through testing first"),
    ("Dropping test data", "Subsequent migrations may depend on it"),
]

for i, (title, desc) in enumerate(prohibitions):
    y = Inches(1.75 + i * 0.68)
    card = add_shape(slide, Inches(6.8), y, Inches(6), Inches(0.6), SOFT_RED, border_color=RED, border_width=Pt(0.5))
    add_text_box(slide, Inches(7.0), y + Inches(0.05), Inches(5.6), Inches(0.3),
                title, font_size=11, color=RED, bold=True)
    add_text_box(slide, Inches(7.0), y + Inches(0.3), Inches(5.6), Inches(0.3),
                desc, font_size=9.5, color=MEDIUM_TEXT)

add_footer(slide)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Multi-Agent Migration Workflow
# ═══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_title_bar(slide, "Multi-Agent Migration Workflow", "Professional agents of a team collaborating on large-scale migration")

# Manager Agent at center
manager_card = add_shape(slide, Inches(5.5), Inches(4.5), Inches(2.4), Inches(1.6), SOFT_BLUE, border_color=DARK_BLUE, border_width=Pt(2.5))
add_text_box(slide, Inches(5.5), Inches(4.55), Inches(2.4), Inches(0.3),
             "Manager Agent", font_size=12, color=DARK_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
manager_desc = [
    "Controls workflow",
    "Dispatches tasks",
    "Coordinates testing",
    "Appends to target"
]
add_text_box(slide, Inches(5.6), Inches(4.9), Inches(2.2), Inches(1.12),
             '\n'.join([f"• {desc}" for desc in manager_desc]), font_size=10, color=DARK_TEXT)

# Requester Agent (top-left)
requester_card = add_shape(slide, Inches(5.5), Inches(1.3), Inches(2.4), Inches(1.6), SOFT_PURPLE, border_color=PURPLE, border_width=Pt(2))
add_text_box(slide, Inches(5.5), Inches(1.35), Inches(2.4), Inches(0.3),
             "Requester Agent", font_size=12, color=PURPLE, bold=True, alignment=PP_ALIGN.CENTER)
requester_desc = [
    "Reads files section-by-section",
    "Uses Read(offset, limit)",
    "Identifies objects",
    "Maintains reading state"
]
add_text_box(slide, Inches(5.6), Inches(1.7), Inches(2.2), Inches(1.12),
             '\n'.join([f"• {desc}" for desc in requester_desc]), font_size=10, color=DARK_TEXT)

# Migrator Agent (bottom-left)
migrator_card = add_shape(slide, Inches(0.5), Inches(4.5), Inches(2.4), Inches(1.6), SOFT_ORANGE, border_color=ORANGE, border_width=Pt(2))
add_text_box(slide, Inches(0.5), Inches(4.55), Inches(2.4), Inches(0.3),
             "Migrator Agent", font_size=12, color=ORANGE, bold=True, alignment=PP_ALIGN.CENTER)
migrator_desc = [
    "Code transformation",
    "Loads docs on-demand",
    "Unit tests code",
    "Returns verified code"
]
add_text_box(slide, Inches(0.6), Inches(4.9), Inches(2.2), Inches(1.12),
             '\n'.join([f"• {desc}" for desc in migrator_desc]), font_size=10, color=DARK_TEXT)

# Tester Agent (right)
tester_card = add_shape(slide, Inches(10.4), Inches(4.5), Inches(2.4), Inches(1.6), SOFT_GREEN, border_color=GREEN, border_width=Pt(2))
add_text_box(slide, Inches(10.4), Inches(4.55), Inches(2.4), Inches(0.3),
             "Tester Agent", font_size=12, color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)
tester_desc = [
    "Functional testing",
    "Integration testing",
    "Complete logs check",
    "Pass/fail feedback"
]
add_text_box(slide, Inches(10.5), Inches(4.9), Inches(2.2), Inches(1.22),
             '\n'.join([f"• {desc}" for desc in tester_desc]), font_size=10, color=DARK_TEXT)

# Arrows between Manager and other agents

# Requester → Manager (REQUEST_READ - arrow pointing down)
arrow_req_mgr = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(6.7), Inches(3.2), Inches(0.3), Inches(1.0))
arrow_req_mgr.fill.solid()
arrow_req_mgr.fill.fore_color.rgb = PURPLE
arrow_req_mgr.line.fill.background()
add_text_box(slide, Inches(6.8), Inches(3.65), Inches(1.7), Inches(0.2), "REQUEST_READ", font_size=6.5, color=PURPLE, bold=True, alignment=PP_ALIGN.CENTER)

# Manager → Requester (READ_RESPONSE - arrow pointing up)
arrow_mgr_req = slide.shapes.add_shape(MSO_SHAPE.UP_ARROW, Inches(6.3), Inches(3.2), Inches(0.3), Inches(1.0))
arrow_mgr_req.fill.solid()
arrow_mgr_req.fill.fore_color.rgb = PURPLE
arrow_mgr_req.line.fill.background()
add_text_box(slide, Inches(4.8), Inches(3.65), Inches(1.7), Inches(0.2), "READ_RESPONSE", font_size=6.5, color=PURPLE, bold=True, alignment=PP_ALIGN.CENTER)

# Migrator → Manager (MIGRATE_REQUEST - arrow pointing right)
arrow_mig_mgr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(3.1), Inches(5.1), Inches(2.3), Inches(0.3))
arrow_mig_mgr.fill.solid()
arrow_mig_mgr.fill.fore_color.rgb = ORANGE
arrow_mig_mgr.line.fill.background()
add_text_box(slide, Inches(3.3), Inches(4.85), Inches(1.7), Inches(0.2), "MIGRATE_REQUEST", font_size=6.5, color=ORANGE, bold=True, alignment=PP_ALIGN.CENTER)

# Manager → Migrator (MIGRATE_RESPONSE - arrow pointing left)
arrow_mgr_mig = slide.shapes.add_shape(MSO_SHAPE.LEFT_ARROW, Inches(3.1), Inches(5.5), Inches(2.3), Inches(0.3))
arrow_mgr_mig.fill.solid()
arrow_mgr_mig.fill.fore_color.rgb = ORANGE
arrow_mgr_mig.line.fill.background()
add_text_box(slide, Inches(3.3), Inches(5.85), Inches(1.7), Inches(0.2), "MIGRATE_RESPONSE", font_size=6.5, color=ORANGE, bold=True, alignment=PP_ALIGN.CENTER)

# Manager → Tester (TEST_REQUEST - arrow pointing right)
arrow_mgr_tst = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(8.0), Inches(5.0), Inches(2.3), Inches(0.3))
arrow_mgr_tst.fill.solid()
arrow_mgr_tst.fill.fore_color.rgb = GREEN
arrow_mgr_tst.line.fill.background()
add_text_box(slide, Inches(8.5), Inches(4.85), Inches(1.7), Inches(0.2), "TEST_REQUEST", font_size=6.5, color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)

# Tester → Manager (TEST_RESPONSE - arrow pointing left)
arrow_tst_mgr = slide.shapes.add_shape(MSO_SHAPE.LEFT_ARROW, Inches(8.0), Inches(5.5), Inches(2.3), Inches(0.3))
arrow_tst_mgr.fill.solid()
arrow_tst_mgr.fill.fore_color.rgb = GREEN
arrow_tst_mgr.line.fill.background()
add_text_box(slide, Inches(8.5), Inches(5.85), Inches(1.7), Inches(0.2), "TEST_RESPONSE", font_size=6.5, color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)

# Benefits section
benefits_box = add_shape(slide, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.6), LIGHT_GRAY, border_color=DARK_BLUE, border_width=Pt(1))
add_text_box(slide, Inches(0.7), Inches(6.32), Inches(12), Inches(0.55),
             "Benefits: Context isolation prevents overflow • Specialized agents ensure rule compliance • Scalable for large migrations • Clear separation of concerns",
             font_size=12, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)

add_footer(slide)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Core Reference Documentation
# ═══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_title_bar(slide, "Core Reference Documentation", "18 comprehensive guides organized in a hierarchical structure")

add_shape(slide, Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.6), SOFT_RED, border_color=RED, border_width=Pt(1.5))
add_text_box(slide, Inches(0.7), Inches(1.35), Inches(12), Inches(0.5),
             "🚨 MANDATORY — Generic Migration Guide: Complete migration requirements that apply to ALL database types (read first!)",
             font_size=13, color=RED, bold=True)

add_shape(slide, Inches(0.5), Inches(2.0), Inches(12.3), Inches(0.6), SOFT_ORANGE, border_color=ORANGE, border_width=Pt(1.5))
add_text_box(slide, Inches(0.7), Inches(2.05), Inches(12), Inches(0.5),
             "🔄 ESSENTIAL — OLTP to OLAP Rewrite Guide: 5 rewrite patterns for converting procedural/row-by-row code to set-based SQL",
             font_size=13, color=ORANGE, bold=True)

guides = [
    ("SQL Syntax Reference", "Complete Vertica SQL syntax: DDL, DML, queries, CTEs, window functions", DARK_BLUE),
    ("Data Types", "Data type mapping & optimization across all source databases", DARK_BLUE),
    ("Function Mapping Guide", "100+ function conversions: Oracle, DB2, SQL Server, PostgreSQL, MySQL", DARK_BLUE),
    ("User-Defined SQL Functions", "CREATE FUNCTION syntax, examples, overloading, and management", DARK_BLUE),
    ("Stored Procedures Guide", "PL/vSQL development: parameters, control structures, exception handling", DARK_BLUE),
    ("UDx Development Guide", "Custom functions in C++, Python, Java, R for scalar, aggregate, analytic, transform", DARK_BLUE),
    ("Query Optimization", "Projection design, encoding strategies, join optimization, statistics", DARK_BLUE),
    ("Machine Learning Guide", "In-database ML: regression, classification, clustering, time series", DARK_BLUE),
    ("ML Function Mapping", "Cross-database ML function equivalents: Python/R ↔ Vertica SQL", DARK_BLUE),
]

for i, (title, desc, color) in enumerate(guides):
    col = i % 3
    row = i // 3
    x = Inches(0.5 + col * 4.15)
    y = Inches(2.8 + row * 1.35)
    card = add_shape(slide, x, y, Inches(4.0), Inches(1.2), LIGHT_GRAY, border_color=color, border_width=Pt(1))
    add_text_box(slide, x + Inches(0.15), y + Inches(0.08), Inches(3.7), Inches(0.35),
                title, font_size=12, color=color, bold=True)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.4), Inches(3.7), Inches(0.7),
                desc, font_size=9.5, color=MEDIUM_TEXT)

add_footer(slide)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — OLTP→OLAP Pattern 1: Adjacent DML Merging
# ═══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_title_bar(slide, "OLTP → OLAP Pattern 1: Adjacent DML → Bulk Operation", "Merge multiple single-row INSERT/UPDATE into one bulk operation")

# Left: Anti-pattern
add_text_box(slide, Inches(0.5), Inches(1.3), Inches(6), Inches(0.35),
             "❌ Anti-Pattern: Multiple Single-Row INSERTs", font_size=14, color=RED, bold=True)
add_code_block(slide, Inches(0.5), Inches(1.65), Inches(6), Inches(1.5),
               """INSERT INTO sales_summary (region, total_sales)
  VALUES ('East', 15000);
INSERT INTO sales_summary (region, total_sales)
  VALUES ('West', 23000);
INSERT INTO sales_summary (region, total_sales)
  VALUES ('North', 18000);
INSERT INTO sales_summary (region, total_sales)
  VALUES ('South', 12000);""")

add_text_box(slide, Inches(0.5), Inches(3.25), Inches(6), Inches(0.35),
             "❌ Anti-Pattern: Multiple Single-Row UPDATEs", font_size=14, color=RED, bold=True)
add_code_block(slide, Inches(0.5), Inches(3.6), Inches(6), Inches(2.1),
               """UPDATE products SET price = price * 1.10
  WHERE product_id = 101;
UPDATE products SET price = price * 1.10
  WHERE product_id = 102;
UPDATE products SET price = price * 1.15
  WHERE product_id = 103;
UPDATE products SET price = price * 1.05
  WHERE product_id = 104;""")

# Right: Optimized
add_text_box(slide, Inches(6.9), Inches(1.3), Inches(6), Inches(0.35),
             "✅ Optimized: Multi-Row INSERT", font_size=14, color=GREEN, bold=True)
add_code_block(slide, Inches(6.9), Inches(1.65), Inches(6), Inches(1.0),
               """INSERT INTO sales_summary (region, total_sales) VALUES
  ('East',  15000),
  ('West',  23000),
  ('North', 18000),
  ('South', 12000);""",
               text_color=RGBColor(0x98, 0xC3, 0x79))

add_text_box(slide, Inches(6.9), Inches(3.25), Inches(6), Inches(0.35),
             "✅ Optimized: Set-Based UPDATE", font_size=14, color=GREEN, bold=True)
add_code_block(slide, Inches(6.9), Inches(3.6), Inches(6), Inches(2.1),
               """UPDATE products
SET price = price * CASE product_id
    WHEN 101 THEN 1.10
    WHEN 102 THEN 1.10
    WHEN 103 THEN 1.15
    WHEN 104 THEN 1.05
END
WHERE product_id IN (101,102,103,104);

-- Or from a lookup table:
UPDATE products p
SET price = p.price * r.factor
FROM price_adjustments r
WHERE p.product_id = r.product_id;""",
               text_color=RGBColor(0x98, 0xC3, 0x79))

# Why it's better
add_shape(slide, Inches(0.5), Inches(6.0), Inches(12.3), Inches(1.0), SOFT_GREEN, border_color=GREEN, border_width=Pt(1.5))
add_text_box(slide, Inches(0.7), Inches(6.05), Inches(12), Inches(0.35),
             "💡 Why it's better:", font_size=13, color=GREEN, bold=True)
add_text_box(slide, Inches(0.7), Inches(6.4), Inches(12), Inches(0.6),
             "• Single parse, single transaction, single bulk write\n• Fewer ROS containers, less Tuple Mover overhead, better compression\n• Vertica's UPDATE is delete+append at storage level — batching minimizes ROS fragmentation",
             font_size=11, color=DARK_TEXT)

add_footer(slide)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — OLTP→OLAP Pattern 2: Loop-DML → Set-Based SQL
# ═══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_title_bar(slide, "OLTP → OLAP Pattern 2: Loop-DML → Set-Based SQL", "Replace row-by-row loops with single SQL statements")

# Anti-pattern
add_text_box(slide, Inches(0.5), Inches(1.3), Inches(6), Inches(0.35),
             "❌ Anti-Pattern: Date Range Loop", font_size=14, color=RED, bold=True)
add_code_block(slide, Inches(0.5), Inches(1.65), Inches(6), Inches(3.3),
               """-- PL/vSQL: iterate over dates,
-- one SELECT + one INSERT per day
DECLARE
    v_date DATE;
    v_count INTEGER;
BEGIN
    FOR v_date IN QUERY
        SELECT generate_series(
            DATE '2024-01-01',
            DATE '2024-12-31',
            INTERVAL '1 day')
    LOOP
        SELECT COUNT(*) INTO v_count
        FROM orders
        WHERE order_date = v_date;

        PERFORM INSERT INTO daily_stats
        VALUES (v_date, v_count);
    END LOOP;
END;""")

# Optimized
add_text_box(slide, Inches(6.9), Inches(1.3), Inches(6), Inches(0.35),
             "✅ Optimized: Single INSERT...SELECT", font_size=14, color=GREEN, bold=True)
add_code_block(slide, Inches(6.9), Inches(1.65), Inches(6), Inches(3.3),
               """-- Single set-based operation:
-- dates with zero orders included
INSERT INTO daily_stats
SELECT g.d::DATE, COUNT(o.order_id)
FROM (
    SELECT generate_series(
        DATE '2024-01-01',
        DATE '2024-12-31',
        INTERVAL '1 day') AS d
) g
LEFT JOIN orders o
    ON o.order_date = g.d
GROUP BY g.d;""",
               text_color=RGBColor(0x98, 0xC3, 0x79))

# Sub-patterns
add_text_box(slide, Inches(0.5), Inches(5.1), Inches(12.3), Inches(0.3),
             "Pattern 2 Sub-Types", font_size=13, color=DARK_BLUE, bold=True)

sub_types = [
    ("2A", "generate_series Loop", "INSERT...SELECT with generate_series / unnest", ORANGE),
    ("2B", "Lookup Table Loop", "INSERT...SELECT with JOIN", GREEN),
    ("2C", "Temp Table + Loop", "CTE or Derived Table", TEAL),
    ("2D", "Row-by-Row DELETE", "Set-Based DELETE or Partition Drop", PURPLE),
    ("2E", "Row-by-Row UPSERT", "Vertica MERGE", RED),
]

for i, (label, title, desc, color) in enumerate(sub_types):
    x = Inches(0.5 + i * 2.55)
    y = Inches(5.45)
    w = Inches(2.45)
    card = add_shape(slide, x, y, w, Inches(1.2), LIGHT_GRAY, border_color=color, border_width=Pt(1))
    badge = add_shape(slide, x + Inches(0.1), y + Inches(0.1), Inches(0.5), Inches(0.35), color)
    add_text_box(slide, x + Inches(0.1), y + Inches(0.1), Inches(0.5), Inches(0.35), label,
                font_size=10, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.7), y + Inches(0.1), Inches(1.6), Inches(0.35), title,
                font_size=10, color=DARK_BLUE, bold=True)
    add_text_box(slide, x + Inches(0.1), y + Inches(0.5), w - Inches(0.2), Inches(1.4),
                desc, font_size=9, color=MEDIUM_TEXT)

add_footer(slide)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — OLTP→OLAP Pattern 3: Cursor → Window Functions
# ═══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_title_bar(slide, "OLTP → OLAP Pattern 3: Cursor → Window Functions", "Replace cursor-based iteration with window/analytic functions, CTE for sharing variables and complex multi-step calculations")

add_text_box(slide, Inches(0.5), Inches(1.3), Inches(6), Inches(0.35),
             "❌ Anti-Pattern: Cursor-Based Row Processing", font_size=14, color=RED, bold=True)
add_code_block(slide, Inches(0.5), Inches(1.65), Inches(6), Inches(2.8),
               """DECLARE
    CURSOR c IS
        SELECT employee_id, salary, hire_date
        FROM employees ORDER BY hire_date;
    v_prev_salary NUMERIC;
    v_emp_id INTEGER;
    v_salary NUMERIC;
    v_hire_date DATE;
BEGIN
    OPEN c;
    FETCH c INTO v_emp_id, v_salary, v_hire_date;
    WHILE FOUND LOOP
        -- process each row with
        -- access to previous row's data
        v_prev_salary := ...;
        FETCH c INTO v_emp_id, v_salary, v_hire_date;
    END LOOP;
    CLOSE c;
END;""")

add_text_box(slide, Inches(6.9), Inches(1.3), Inches(6), Inches(0.35),
             "✅ Optimized: Window Functions", font_size=14, color=GREEN, bold=True)
add_code_block(slide, Inches(6.9), Inches(1.65), Inches(6), Inches(2.8),
               """SELECT employee_id, salary, hire_date,
       LAG(salary) OVER (
           ORDER BY hire_date
       ) AS prev_salary,
       LEAD(salary) OVER (
           ORDER BY hire_date
       ) AS next_salary,
       ROW_NUMBER() OVER (
           ORDER BY hire_date
       ) AS row_num,
       SUM(salary) OVER (
           ORDER BY hire_date
           ROWS UNBOUNDED PRECEDING
       ) AS running_total
FROM employees;""",
               text_color=RGBColor(0x98, 0xC3, 0x79))

# Mapping table
add_text_box(slide, Inches(0.5), Inches(4.6), Inches(12.3), Inches(0.3),
             "Cursor-to-Window Function Mapping", font_size=13, color=DARK_BLUE, bold=True)

mappings = [
    ("CURSOR + FETCH loop", "Window function in SELECT"),
    ("Access previous row", "LAG(col) OVER (ORDER BY ...)"),
    ("Access next row", "LEAD(col) OVER (ORDER BY ...)"),
    ("Row counter", "ROW_NUMBER() OVER (ORDER BY ...)"),
    ("Running total", "SUM(col) OVER (ORDER BY ... ROWS UNBOUNDED PRECEDING)"),
    ("Group ranking", "RANK() / DENSE_RANK() OVER (PARTITION BY ... ORDER BY ...)"),
]

for i, (cursor, window) in enumerate(mappings):
    x = Inches(0.5 + (i % 3) * 4.15)
    y = Inches(4.95 + (i // 3) * 0.85)
    card = add_shape(slide, x, y, Inches(4.0), Inches(0.75), SOFT_TEAL, border_color=TEAL, border_width=Pt(0.5))
    add_text_box(slide, x + Inches(0.1), y + Inches(0.05), Inches(1.8), Inches(0.3),
                cursor, font_size=9, color=RED, bold=True)
    add_text_box(slide, x + Inches(2.0), y + Inches(0.05), Inches(0.3), Inches(0.3), "→",
                font_size=12, color=DARK_TEXT, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(2.4), y + Inches(0.05), Inches(1.5), Inches(0.3),
                window, font_size=9, color=GREEN, bold=True)
    add_text_box(slide, x + Inches(0.05), y + Inches(0.35), Inches(2.6), Inches(0.35),
                f"Replace {cursor.lower()} with {window.lower()}", font_size=8.5, color=MEDIUM_TEXT)

add_footer(slide)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — OLTP→OLAP Pattern 4 & 5
# ═══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_title_bar(slide, "OLTP → OLAP Patterns 4 & 5", "Function-call→JOIN and Recursive CTE")

# Pattern 4
add_text_box(slide, Inches(0.5), Inches(1.3), Inches(6), Inches(0.35),
             "❌ Pattern 4: Per-Row Function Call", font_size=14, color=RED, bold=True)
add_code_block(slide, Inches(0.5), Inches(1.65), Inches(6), Inches(1.8),
               """-- Oracle: function called per row
SELECT employee_id,
       get_dept_name(dept_id)
       AS dept_name,
       get_manager_name(mgr_id)
       AS manager_name
FROM employees;

-- Each function call is a
-- separate query per row!
-- N rows × M functions = N×M queries""")

add_text_box(slide, Inches(6.9), Inches(1.3), Inches(6), Inches(0.35),
             "✅ Optimized: JOIN to Lookup Tables", font_size=14, color=GREEN, bold=True)
add_code_block(slide, Inches(6.9), Inches(1.65), Inches(6), Inches(1.8),
               """-- Single query with JOINs
SELECT e.employee_id,
       d.dept_name,
       m.name AS manager_name
FROM employees e
LEFT JOIN departments d
    ON e.dept_id = d.dept_id
LEFT JOIN employees m
    ON e.mgr_id = m.employee_id;""",
               text_color=RGBColor(0x98, 0xC3, 0x79))

# Pattern 5
add_text_box(slide, Inches(0.5), Inches(3.6), Inches(6), Inches(0.35),
             "❌ Pattern 5: Oracle CONNECT BY", font_size=14, color=RED, bold=True)
add_code_block(slide, Inches(0.5), Inches(3.95), Inches(6), Inches(2.1),
               """-- Oracle: hierarchical query
SELECT employee_id, manager_id,
       name, LEVEL,
       SYS_CONNECT_BY_PATH(name,'/')
       AS path
FROM employees
START WITH manager_id IS NULL
CONNECT BY PRIOR employee_id
    = manager_id;""")

add_text_box(slide, Inches(6.9), Inches(3.6), Inches(6), Inches(0.35),
             "✅ Optimized: WITH RECURSIVE CTE", font_size=14, color=GREEN, bold=True)
add_code_block(slide, Inches(6.9), Inches(3.95), Inches(6), Inches(2.1),
               """-- Vertica: recursive CTE
WITH RECURSIVE emp_tree AS (
    SELECT employee_id, manager_id,
           name, 1 AS level,
           '/' || name AS path
    FROM employees
    WHERE manager_id IS NULL
    UNION ALL
    SELECT e.employee_id, e.manager_id,
           e.name, t.level + 1,
           t.path || '/' || e.name
    FROM employees e
    JOIN emp_tree t
        ON e.manager_id = t.employee_id
)
SELECT * FROM emp_tree;""",
               text_color=RGBColor(0x98, 0xC3, 0x79))

# Syntax mapping
add_shape(slide, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.7), SOFT_PURPLE, border_color=PURPLE, border_width=Pt(1))
add_text_box(slide, Inches(0.7), Inches(6.35), Inches(11.9), Inches(0.3),
             "CONNECT BY → WITH RECURSIVE Syntax Mapping", font_size=12, color=PURPLE, bold=True)
add_text_box(slide, Inches(0.7), Inches(6.55), Inches(11.9), Inches(0.4),
             "START WITH → Anchor WHERE  |  CONNECT BY PRIOR → Recursive JOIN  |  LEVEL → Manual counter  |  SYS_CONNECT_BY_PATH → String concat  |  NOCYCLE → Manual depth limit",
             font_size=10.5, color=DARK_TEXT)

add_footer(slide)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — User-Defined SQL Functions
# ═══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_title_bar(slide, "User-Defined SQL Functions", "Extend Vertica with custom SQL expressions")

add_text_box(slide, Inches(0.5), Inches(1.3), Inches(5.8), Inches(0.4),
             "CREATE FUNCTION Syntax", font_size=16, color=DARK_BLUE, bold=True)

add_code_block(slide, Inches(0.5), Inches(1.7), Inches(5.8), Inches(2.0),
               """CREATE [OR REPLACE] FUNCTION
    schema.func_name(
        arg1 INT,
        arg2 VARCHAR(50)
    )
RETURN INT
AS BEGIN
    RETURN (
        CASE WHEN arg1 IS NOT NULL
             THEN arg1
             ELSE 0
        END
    );
END;""")

add_text_box(slide, Inches(6.8), Inches(1.3), Inches(6), Inches(0.4),
             "When to Use", font_size=16, color=DARK_BLUE, bold=True)

yes_items = [
    "Simple SQL expressions & calculations",
    "Data transformation & cleaning",
    "Business rule encapsulation",
    "Single RETURN statement logic",
]
no_items = [
    "Complex procedural logic with loops",
    "Functions needing FROM/WHERE/GROUP BY",
    "Aggregate or analytic functions",
    "Complex data types (ARRAY, ROW, SET)",
]

add_text_box(slide, Inches(6.8), Inches(1.75), Inches(2.8), Inches(0.3), "✅ Perfect for:",
             font_size=12, color=GREEN, bold=True)
add_bullet_list(slide, Inches(6.8), Inches(2.05), Inches(2.8), Inches(1.6), yes_items, font_size=11, color=DARK_TEXT)
add_text_box(slide, Inches(10.0), Inches(1.75), Inches(2.8), Inches(0.3), "❌ Not suitable for:",
             font_size=12, color=RED, bold=True)
add_bullet_list(slide, Inches(10.0), Inches(2.05), Inches(2.8), Inches(1.6), no_items, font_size=11, color=DARK_TEXT)

add_text_box(slide, Inches(0.5), Inches(3.9), Inches(12.3), Inches(0.4),
             "Example: NULL Handling Function", font_size=16, color=DARK_BLUE, bold=True)

add_code_block(slide, Inches(0.5), Inches(4.3), Inches(12.3), Inches(2.5),
               """-- Create function to replace NULL with zero
CREATE FUNCTION myzeroifnull(x INT) RETURN INT
   AS BEGIN
     RETURN (CASE WHEN (x IS NOT NULL) THEN x ELSE 0 END);
   END;

-- Usage examples
SELECT myzeroifnull(column_name) FROM table_name;
SELECT COUNT(*) FROM table_name
GROUP BY myzeroifnull(column_name);

-- Functions are flattened and optimized by the query planner""")

add_footer(slide)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — PL/vSQL Stored Procedures
# ═══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_title_bar(slide, "PL/vSQL Stored Procedures", "Vertica's procedural language for stored procedures")

add_text_box(slide, Inches(0.5), Inches(1.3), Inches(6), Inches(0.4),
             "Key Characteristics", font_size=16, color=DARK_BLUE, bold=True)

chars = [
    ("OLAP Optimized", "Designed for analytical processing, not OLTP row-by-row operations", ORANGE),
    ("Set-Based", "Set-based operations preferred over row-by-row loops", GREEN),
    ("Exception Handling", "GET STACKED DIAGNOSTICS for comprehensive error management", TEAL),
    ("Nested Support", "Up to 50 levels of nested procedure calls", PURPLE),
    ("PL/pgSQL Compatible", "Based on PostgreSQL's PL/pgSQL with minor differences", RED),
]

for i, (title, desc, color) in enumerate(chars):
    x = Inches(0.5 + (i % 3) * 4.15)
    y = Inches(1.8 + (i // 3) * 1.1)
    card = add_shape(slide, x, y, Inches(4.0), Inches(0.95), LIGHT_GRAY, border_color=color, border_width=Pt(1))
    add_text_box(slide, x + Inches(0.15), y + Inches(0.08), Inches(3.7), Inches(0.3),
                title, font_size=12, color=color, bold=True)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.4), Inches(3.7), Inches(0.5),
                desc, font_size=10, color=MEDIUM_TEXT)

add_text_box(slide, Inches(0.5), Inches(4.0), Inches(12.3), Inches(0.4),
             "PL/vSQL-Only Commands", font_size=16, color=DARK_BLUE, bold=True)

commands = [
    ("PERFORM", "Execute SQL and discard result (DDL, DML, CALL, COMMIT, ROLLBACK)", "PERFORM INSERT INTO audit_log VALUES ('started');"),
    ("RAISE", "Error handling and messaging (NOTICE, WARNING, EXCEPTION)", "RAISE NOTICE 'Processing % rows';"),
    (":= / <-", "Variable assignment (regular / truncating)", "v_count := v_count + 1;"),
    ("FOUND", "Check if last SQL returned rows", "IF FOUND THEN RAISE NOTICE 'Success'; END IF;"),
]

for i, (cmd, desc, example) in enumerate(commands):
    x = Inches(0.5 + (i % 2) * 6.2)
    y = Inches(4.45 + (i // 2) * 1.3)
    card = add_shape(slide, x, y, Inches(6.0), Inches(1.15), SOFT_BLUE, border_color=ACCENT_BLUE, border_width=Pt(1))
    add_text_box(slide, x + Inches(0.15), y + Inches(0.05), Inches(5.7), Inches(0.3),
                cmd, font_size=13, color=DARK_BLUE, bold=True)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.35), Inches(5.7), Inches(0.3),
                desc, font_size=10, color=MEDIUM_TEXT)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.65), Inches(5.7), Inches(0.35),
                example, font_size=9.5, color=RGBColor(0x00, 0x64, 0x00), font_name="Consolas")

add_footer(slide)



# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Parameter DEFAULT Keyword Limitation
# ═══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_title_bar(slide, "Parameter DEFAULT Keyword Limitation & Solution", "Vertica PL/vSQL does NOT support DEFAULT in parameter declarations")

# Left: Problem
add_text_box(slide, Inches(0.5), Inches(1.3), Inches(6), Inches(0.35),
             "❌ Syntax Error in Vertica", font_size=14, color=RED, bold=True)

add_code_block(slide, Inches(0.5), Inches(1.65), Inches(6), Inches(2.0),
               """-- This will NOT work in Vertica
CREATE PROCEDURE process_order(
    p_order_id INTEGER,
    p_discount FLOAT DEFAULT 0.1, --  <- Syntax error!
    p_priority VARCHAR DEFAULT 'NORMAL' -- <- Syntax error!
) AS $$
BEGIN
    RAISE NOTICE 'Processing order: %', p_order_id;
END;
$$;""")

add_shape(slide, Inches(0.5), Inches(3.8), Inches(6), Inches(0.8), SOFT_RED, border_color=RED, border_width=Pt(1))
add_text_box(slide, Inches(0.7), Inches(3.85), Inches(5.6), Inches(0.7),
             "🚨 PL/vSQL does not support the DEFAULT keyword in procedure parameter declarations.",
             font_size=11, color=RED, bold=True)

# Right: Solution
add_text_box(slide, Inches(6.9), Inches(1.3), Inches(6), Inches(0.35),
             "✅ Solution: Procedure Overloading", font_size=14, color=GREEN, bold=True)

add_code_block(slide, Inches(6.9), Inches(1.65), Inches(6), Inches(5.15),
               """-- Main procedure (all parameters)
CREATE PROCEDURE process_order(
    p_order_id INTEGER,
    p_discount FLOAT,
    p_priority VARCHAR
) AS $$
BEGIN
    -- Core business logic here
    RAISE NOTICE 'Processing order: %', p_order_id;
END;
$$;

-- Overloaded: 1 param (2 defaults)
CREATE PROCEDURE process_order(
    p_order_id INTEGER
) AS $$
BEGIN
    PERFORM CALL process_order(
        p_order_id, 0.1, 'NORMAL');
END;
$$;

-- Overloaded: 2 params (1 defaults)
CREATE PROCEDURE process_order(
    p_order_id INTEGER,
    p_discount FLOAT
) AS $$
BEGIN
    PERFORM CALL process_order(
        p_order_id, p_discount, 'NORMAL');
END;
$$;""",
               text_color=RGBColor(0x98, 0xC3, 0x79))

# Key points
add_shape(slide, Inches(0.5), Inches(4.8), Inches(6), Inches(2.0), SOFT_GREEN, border_color=GREEN, border_width=Pt(1.5))
add_text_box(slide, Inches(0.7), Inches(4.85), Inches(6), Inches(0.35),
             "💡 Key Points:", font_size=13, color=GREEN, bold=True)
add_text_box(slide, Inches(0.7), Inches(5.2), Inches(6), Inches(1.5),
             "• All overloaded procedures MUST have the EXACT SAME NAME \u2014 only the parameter list differs\n"
             "• Business logic exists ONLY in the main procedure \u2014 overloaded versions just call it with defaults\n"
             "• Calling patterns are 100% compatible with Oracle and other databases", 
             font_size=11, color=DARK_TEXT)

add_footer(slide)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Function Migration Strategies (Oracle & DB2)
# ═══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_title_bar(slide, "Function Migration Strategies", "Different approaches for different scenarios")

# Strategy 1
add_text_box(slide, Inches(0.5), Inches(1.3), Inches(6), Inches(0.35),
             "Strategy 1: SQL Function → Subquery (Performance)", font_size=14, color=ORANGE, bold=True)
add_code_block(slide, Inches(0.5), Inches(1.65), Inches(6), Inches(2.2),
               """-- Oracle/DB2: function queries table
CREATE FUNCTION ISYSZ(rydm VARCHAR)
RETURNS VARCHAR
BEGIN
  DECLARE rynum INT;
  SELECT COUNT(*) INTO rynum
  FROM qx_user WHERE czry_dm = rydm;
  IF rynum > 0 THEN RETURN '1';
  ELSE RETURN '0'; END IF;
END;

-- Called per row: N row scans
SELECT czry_dm, ISYSZ(czry_dm)
FROM dm_czry;""")

add_code_block(slide, Inches(6.9), Inches(1.65), Inches(6), Inches(2.2),
               """-- Vertica: LEFT JOIN subquery
SELECT dm.czry_dm, dm.czry_mc,
  CASE WHEN u.userid IS NOT NULL
       THEN '1' ELSE '0'
  END AS isysz
FROM dm_czry dm
LEFT JOIN qx_user u
    ON dm.czry_dm = u.czry_dm;

-- ✅ Single scan, set-based""",
               text_color=RGBColor(0x98, 0xC3, 0x79))

# Strategy 2
add_text_box(slide, Inches(0.5), Inches(4.0), Inches(6), Inches(0.35),
             "Strategy 2: Function → Stored Procedure (Complex Logic)", font_size=14, color=GREEN, bold=True)
add_code_block(slide, Inches(0.5), Inches(4.35), Inches(6), Inches(2.0),
               """-- Oracle/DB2 function
CREATE FUNCTION F_GET_JDH()
RETURNS VARCHAR
BEGIN
  DECLARE jdno VARCHAR;
  SELECT CSNR INTO jdno
  FROM XT_XTCS
  WHERE CSXH = '10001';
  RETURN jdno;
END;

-- Usage: jdno := F_GET_JDH();""")

add_code_block(slide, Inches(6.9), Inches(4.35), Inches(6), Inches(2.0),
               """-- Vertica stored procedure
CREATE PROCEDURE F_GET_JDH(
  OUT rt VARCHAR(100)
) AS $$
BEGIN
  SELECT CSNR INTO rt
  FROM XT_XTCS
  WHERE CSXH = '10001';
END;
$$;

-- Tuple unpacking assignment
jdno := CALL F_GET_JDH();""",
               text_color=RGBColor(0x98, 0xC3, 0x79))

# Decision table
add_text_box(slide, Inches(0.7), Inches(6.4), Inches(4), Inches(0.3), "Decision Guide:", font_size=11, color=TEAL, bold=True)
add_shape(slide, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.4), SOFT_TEAL, border_color=TEAL, border_width=Pt(1))
add_text_box(slide, Inches(0.7), Inches(6.75), Inches(12), Inches(0.25),
             "Table lookup / SELECT clause → Subquery (Strategy 1)  |  Complex logic / multi-statement → Stored Procedure (Strategy 2)  |  Simple math → User-Defined SQL Function",
             font_size=9.5, color=DARK_TEXT)

add_footer(slide)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — Common Migration Challenges (Oracle & DB2)
# ═══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_title_bar(slide, "Common Migration Challenges", "Oracle & DB2 — Key differences and solutions")

challenges = [
    ("1", "Identifier Case Sensitivity",
     "Oracle/DB2 quoted identifiers are case-sensitive.\nVertica identifiers are always case-insensitive.",
     '"MyTable" vs "mytable" conflict in Vertica → Rename to snake_case',
     ORANGE),
    ("2", "ON DELETE CASCADE",
     "Vertica does NOT support ON DELETE CASCADE\nfor foreign key constraints.",
     "Use stored procedures or application logic for cascade deletes",
     GREEN),
    ("3", "PIVOT / UNPIVOT",
     "Vertica does NOT support PIVOT or UNPIVOT operators.",
     "Rewrite as CASE + aggregate (PIVOT) or UNION ALL (UNPIVOT)",
     TEAL),
    ("4", "CONNECT BY (Oracle)",
     "Oracle's CONNECT BY has no direct equivalent in Vertica.",
     "Rewrite as WITH RECURSIVE CTE with anchor + recursive terms",
     PURPLE),
]

for i, (num, title, problem, solution, color) in enumerate(challenges):
    x = Inches(0.5 + (i % 2) * 6.2)
    y = Inches(1.35 + (i // 2) * 2.7)
    card = add_shape(slide, x, y, Inches(6.0), Inches(2.55), LIGHT_GRAY, border_color=color, border_width=Pt(1.5))
    # Badge
    badge = add_shape(slide, x + Inches(0.15), y + Inches(0.1), Inches(0.45), Inches(0.45), color)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.1), Inches(0.45), Inches(0.45), num,
                font_size=12, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.7), y + Inches(0.1), Inches(5.1), Inches(0.4),
                title, font_size=13, color=color, bold=True)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.6), Inches(5.7), Inches(0.7),
                problem, font_size=10, color=RED)
    add_text_box(slide, x + Inches(0.15), y + Inches(1.2), Inches(5.7), Inches(0.5),
                f"✅ {solution}", font_size=10, color=GREEN, bold=True)
    # Code example for each
    if i == 0:
        add_code_block(slide, x + Inches(0.15), y + Inches(1.5), Inches(5.7), Inches(0.9),
                       '-- "MyTable" vs "mytable" conflict\nCREATE TABLE my_table (id INT);  -- renamed',
                       font_size=8)
    elif i == 1:
        add_code_block(slide, x + Inches(0.15), y + Inches(1.5), Inches(5.7), Inches(0.9),
                       '-- ON DELETE CASCADE removed\nFOREIGN KEY (module_code)\n  REFERENCES modules(module_code)\n  -- ON DELETE CASCADE (not supported)',
                       font_size=8)
    elif i == 2:
        add_code_block(slide, x + Inches(0.15), y + Inches(1.5), Inches(5.7), Inches(0.9),
                       "-- PIVOT → CASE + GROUP BY\nSUM(CASE WHEN q='Q1' THEN sales END) AS Q1",
                       font_size=8)
    elif i == 3:
        add_code_block(slide, x + Inches(0.15), y + Inches(1.5), Inches(5.7), Inches(0.9),
                       "-- CONNECT BY → WITH RECURSIVE\nWITH RECURSIVE cte AS (\n  SELECT ... WHERE manager_id IS NULL\n  UNION ALL\n  SELECT ... JOIN cte ...\n)",
                       font_size=8)

add_footer(slide)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — Common Migration Challenges (SQL Server)
# ═══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_title_bar(slide, "Common Migration Challenges", "SQL Server — Key differences and solutions")

# Challenges grid
ss_challenges = [
    ("IDENTITY Columns", "INT IDENTITY(1,1)\nvs IDENTITY data type", "IDENTITY is standalone type\nDo NOT prefix with INT", ORANGE),
    ("Temporary Tables", "#temp / ##global_temp\nvs LOCAL/GLOBAL TEMP", "Data always session-private\nON COMMIT PRESERVE ROWS", GREEN),
    ("Dynamic SQL", "EXEC(@SQL)\nsp_executesql", "EXECUTE 'SQL' USING params\nIn PL/vSQL only", TEAL),
    ("Cursors", "DECLARE CURSOR\nFETCH NEXT", "Prefer window functions\nor set-based SQL", PURPLE),
    ("JSON Support", "JSON_VALUE()\nOPENJSON()", "Flex Tables with VMap\ncompute_flextable_keys()", RED),
    ("Computed Columns", "AS expression\nPERSISTED", "DEFAULT USING\nREFRESH_COLUMNS()", DARK_BLUE),
]

for i, (title, source, vertica, color) in enumerate(ss_challenges):
    col = i % 3
    row = i // 3
    x = Inches(0.5 + col * 4.15)
    y = Inches(1.35 + row * 2.7)
    card = add_shape(slide, x, y, Inches(4.0), Inches(2.55), LIGHT_GRAY, border_color=color, border_width=Pt(1.5))
    add_text_box(slide, x + Inches(0.15), y + Inches(0.1), Inches(3.7), Inches(0.35),
                title, font_size=13, color=color, bold=True)
    # Source
    add_text_box(slide, x + Inches(0.15), y + Inches(0.5), Inches(1.8), Inches(0.3),
                "SQL Server:", font_size=9, color=RED, bold=True)
    add_code_block(slide, x + Inches(0.1), y + Inches(0.75), Inches(1.9), Inches(1.0),
                   source, font_size=9, bg_color=SOFT_RED, text_color=RGBColor(0x1, 0x1, 0x01))
    # Arrow
    add_text_box(slide, x + Inches(2.0), y + Inches(1.0), Inches(0.3), Inches(0.4), "→",
                font_size=18, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    # Vertica
    add_text_box(slide, x + Inches(2.35), y + Inches(0.5), Inches(1.5), Inches(0.3),
                "Vertica:", font_size=9, color=GREEN, bold=True)
    add_code_block(slide, x + Inches(2.3), y + Inches(0.75), Inches(1.6), Inches(1.0),
                   vertica, font_size=9, bg_color=SOFT_GREEN, text_color=RGBColor(0x1, 0x1, 0x01))

# Key notes
add_shape(slide, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.9), SOFT_ORANGE, border_color=ORANGE, border_width=Pt(1))
add_text_box(slide, Inches(0.7), Inches(6.2), Inches(3), Inches(0.3), "⚠️ Key Notes:", font_size=11, color=ORANGE, bold=True)
add_text_box(slide, Inches(0.7), Inches(6.45), Inches(12), Inches(0.6),
             "IDENTITY is a standalone type — do NOT write INT IDENTITY  |  Temp table data is ALWAYS session-private (no ##global_temp equivalent)\nSELECT INTO TEMP requires ON COMMIT PRESERVE ROWS or data is lost  |  Vertica has no native JSON type — use Flex Tables",
             font_size=9.5, color=DARK_TEXT)

add_footer(slide)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 18 — UDx Development (HIDDEN)
# ═══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.Hidden = True  # Hide this slide
add_background(slide, WHITE)
add_title_bar(slide, "UDx Custom Function Development", "User-Defined Extensions in C++, Python, Java, and R")

udx_types = [
    ("UDSF", "Scalar Functions", "Single row input → single value output\nCan replace any native function", "C++, Python,\nJava, R", ORANGE, SOFT_ORANGE),
    ("UDAF", "Aggregate Functions", "Process one column → one output column\nSupport distributed aggregation", "C++\n(performance)", GREEN, SOFT_GREEN),
    ("UDAnF", "Analytic Functions", "Read multiple input rows\nSupport OVER() clause", "C++,\nJava", TEAL, SOFT_TEAL),
    ("UDTF", "Transform Functions", "Operate on table partitions\nReturn zero or more rows", "C++, Python,\nJava, R", PURPLE, SOFT_PURPLE),
    ("UDL", "Load Functions", "Custom data loading\nSources, filters, parsers", "C++, Java,\nPython", RED, SOFT_RED),
]

for i, (abbr, title, desc, langs, color, bg) in enumerate(udx_types):
    x = Inches(0.5 + i * 2.55)
    y = Inches(1.4)
    w = Inches(2.4)
    card = add_shape(slide, x, y, w, Inches(5.2), bg, border_color=color, border_width=Pt(2))
    badge = add_shape(slide, x + Inches(0.15), y + Inches(0.2), Inches(0.9), Inches(0.5), color)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.22), Inches(0.9), Inches(0.45), abbr,
                font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(1.0), y + Inches(0.22), Inches(1.3), Inches(0.45), title,
                font_size=11, color=DARK_BLUE, bold=True)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.9), w - Inches(0.3), Inches(1.8),
                desc, font_size=10, color=DARK_TEXT)
    add_text_box(slide, x + Inches(0.15), y + Inches(4.5), w - Inches(0.3), Inches(0.4),
                f"Languages:\n{langs}", font_size=10, color=color, bold=True)

add_footer(slide)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 19 — Machine Learning in Vertica
# ═══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.Hidden = True  # Hidden slide
add_background(slide, WHITE)
add_title_bar(slide, "In-Database Machine Learning", "Complete ML lifecycle without data movement")

categories = [
    ("Regression", "Predict continuous numerical values", ["Linear Regression", "XGBoost Regression", "Random Forest Regression", "SVM Regression", "Poisson Regression"], ORANGE, SOFT_ORANGE),
    ("Classification", "Predict categorical outcomes", ["Logistic Regression", "XGBoost Classifier", "Random Forest Classifier", "Naive Bayes", "SVM Classifier"], GREEN, SOFT_GREEN),
    ("Clustering", "Group similar data points", ["K-Means", "Bisecting K-Means", "K-Prototypes"], TEAL, SOFT_TEAL),
    ("Time Series", "Forecast temporal patterns", ["Autoregression", "Moving Average", "ARIMA"], PURPLE, SOFT_PURPLE),
]

for i, (cat, desc, algorithms, color, bg) in enumerate(categories):
    x = Inches(0.5 + i * 3.12)
    y = Inches(1.4)
    w = Inches(3.0)
    card = add_shape(slide, x, y, w, Inches(2.8), bg, border_color=color, border_width=Pt(2))
    add_text_box(slide, x + Inches(0.15), y + Inches(0.2), w - Inches(0.3), Inches(0.35),
                cat, font_size=15, color=color, bold=True)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.55), w - Inches(0.3), Inches(0.3),
                desc, font_size=9.5, color=MEDIUM_TEXT)
    algo_text = "\n".join(f"  • {a}" for a in algorithms)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.9), w - Inches(0.3), Inches(1.8),
                algo_text, font_size=10.5, color=DARK_TEXT)

add_text_box(slide, Inches(0.5), Inches(4.4), Inches(12.3), Inches(0.4),
             "ML Workflow Example: Customer Churn Prediction", font_size=16, color=DARK_BLUE, bold=True)

add_code_block(slide, Inches(0.5), Inches(4.8), Inches(12.3), Inches(2.1),
               """-- 1. Train churn prediction model
SELECT RF_CLASSIFIER('churn_model', 'customer_data',
    'churned', 'age, income, tenure, support_calls');

-- 2. Make real-time predictions
SELECT customer_id,
       PREDICT_RF_CLASSIFIER(age, income, tenure, support_calls
           USING PARAMETERS model_name='churn_model')
       AS churn_probability
FROM new_customers;

-- 3. Evaluate model performance
SELECT RF_EVALUATOR('churn_model', 'test_data',
    'churned', 'age, income, tenure, support_calls');""")

add_footer(slide)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 20 — Performance Optimization
# ═══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.Hidden = True  # Hidden slide
add_background(slide, WHITE)
add_title_bar(slide, "Performance Optimization", "Leveraging Vertica's columnar MPP architecture")

areas = [
    ("Projection Design", [
        ("Order-Optimized", "Sort projections by frequently filtered columns"),
        ("Aggregate Projections", "Pre-compute aggregations for summary queries"),
        ("Replicated Projections", "Small lookup tables on all nodes"),
        ("Hash Segmentation", "Even data distribution across nodes"),
    ], ORANGE, SOFT_ORANGE),
    ("Encoding Strategies", [
        ("RLE", "Run-Length Encoding for low-cardinality sorted columns"),
        ("DELTA", "Delta encoding for sequential numeric data"),
        ("GZIP", "High compression for high-cardinality columns"),
        ("LZO", "Fast compression for intermediate cardinality"),
    ], GREEN, SOFT_GREEN),
    ("Query Optimization", [
        ("ANALYZE_STATISTICS()", "Update statistics after data loads"),
        ("EXPLAIN", "Review query plans for bottlenecks"),
        ("Directed Queries", "Force specific projections for critical queries"),
        ("Workload Analyzer", "Get automated optimization recommendations"),
    ], TEAL, SOFT_TEAL),
    ("Resource Management", [
        ("Resource Pools", "Allocate memory for different workload types"),
        ("Priority Scheduling", "Set query execution priorities"),
        ("Query Monitoring", "Track performance via system tables"),
        ("Tuple Mover", "Optimize mergeout operations"),
    ], PURPLE, SOFT_PURPLE),
]

for i, (area, items, color, bg) in enumerate(areas):
    x = Inches(0.5 + i * 3.12)
    y = Inches(1.4)
    w = Inches(3.0)
    card = add_shape(slide, x, y, w, Inches(5.2), bg, border_color=color, border_width=Pt(2))
    add_text_box(slide, x + Inches(0.15), y + Inches(0.2), w - Inches(0.3), Inches(0.35),
                area, font_size=14, color=color, bold=True)
    for j, (title, desc) in enumerate(items):
        iy = y + Inches(0.7 + j * 1.05)
        add_text_box(slide, x + Inches(0.15), iy, Inches(2.7), Inches(0.28),
                    title, font_size=11, color=DARK_BLUE, bold=True)
        add_text_box(slide, x + Inches(0.15), iy + Inches(0.26), Inches(2.7), Inches(0.7),
                    desc, font_size=9, color=MEDIUM_TEXT)

add_footer(slide)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 21 — VSQL Testing Framework
# ═══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_title_bar(slide, "VSQL Testing Framework", "Immediately test all converted SQL and stored procedures")

add_text_box(slide, Inches(0.5), Inches(1.3), Inches(5.8), Inches(0.4),
             "VSQL Setup", font_size=16, color=DARK_BLUE, bold=True)

add_code_block(slide, Inches(0.5), Inches(1.7), Inches(5.8), Inches(1.2),
               """export VSQL='/opt/vertica/bin/vsql \\
  -h hostname -p 5433 \\
  -U username -w password dbname'

$VSQL -c "SET SESSION AUTOCOMMIT TO ON;\"""")

add_text_box(slide, Inches(6.8), Inches(1.3), Inches(6), Inches(0.4),
             "Check Object Availability", font_size=16, color=DARK_BLUE, bold=True)

add_code_block(slide, Inches(6.8), Inches(1.7), Inches(6), Inches(2.0),
               """Schema:    $VSQL -c "\\dn schema_name"
Table:     $VSQL -c "\\dt table_name"
View:      $VSQL -c "\\dt view_name"
Projection: $VSQL -c "\\dj proj_name"
Function:  $VSQL -c "\\df func_name\\""""",
               text_color=RGBColor(0x98, 0xC3, 0x79))

add_text_box(slide, Inches(0.5), Inches(3.0), Inches(12.3), Inches(0.4),
             "Testing Methods", font_size=16, color=DARK_BLUE, bold=True)

m1_box = add_shape(slide, Inches(0.5), Inches(3.4), Inches(5.8), Inches(1.6), SOFT_BLUE, border_color=ACCENT_BLUE, border_width=Pt(1))
add_text_box(slide, Inches(0.7), Inches(3.45), Inches(5.4), Inches(0.3),
             "Method 1: Individual Commands", font_size=12, color=DARK_BLUE, bold=True)
add_code_block(slide, Inches(0.7), Inches(3.75), Inches(5.4), Inches(1.1),
               """$VSQL -c "CREATE TABLE test (id INT);"
$VSQL -c "INSERT INTO test VALUES(1);
           COMMIT;"
$VSQL -c "SELECT * FROM test;\\""""",
               font_size=8.5)

m2_box = add_shape(slide, Inches(6.8), Inches(3.4), Inches(6), Inches(1.6), SOFT_GREEN, border_color=GREEN, border_width=Pt(1))
add_text_box(slide, Inches(7.0), Inches(3.45), Inches(5.6), Inches(0.3),
             "Method 2: Multi-Statement (Recommended)", font_size=12, color=DARK_BLUE, bold=True)
add_code_block(slide, Inches(7.0), Inches(3.75), Inches(5.6), Inches(1.1),
               """$VSQL<<-'EOF'
  SET SESSION AUTOCOMMIT TO ON;
  CREATE TABLE test (id INTEGER);
  INSERT INTO test VALUES (1);
  SELECT * FROM test;
EOF""",
               font_size=8.5)

add_text_box(slide, Inches(0.5), Inches(5.2), Inches(12.3), Inches(0.35),
             "Test-First Rule: MIGRATE → TEST → PASS → APPEND", font_size=14, color=DARK_BLUE, bold=True)

test_steps = [
    ("MIGRATE", "Convert to Vertica syntax", ORANGE),
    ("TEST", "Execute immediately", GREEN),
    ("PASS", "Only if succeeds → consult docs → retry", TEAL),
    ("APPEND", "Only after passing", PURPLE),
]

for i, (step, desc, color) in enumerate(test_steps):
    x = Inches(0.5 + i * 3.12)
    y = Inches(5.6)
    card = add_shape(slide, x, y, Inches(3.0), Inches(1.2), LIGHT_GRAY, border_color=color, border_width=Pt(1.5))
    add_text_box(slide, x, y + Inches(0.1), Inches(3.0), Inches(0.4), step,
                font_size=16, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.5), Inches(2.7), Inches(0.6),
                desc, font_size=10, color=MEDIUM_TEXT, alignment=PP_ALIGN.CENTER)

add_footer(slide)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 22 — Summary & Key Takeaways
# ═══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, DARK_BLUE)

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8),
             "Summary & Key Takeaways", font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

takeaways = [
    ("5", "Migration Paths", "Oracle, DB2, SQL Server, PostgreSQL, MySQL → Vertica", ORANGE),
    ("100+", "Function Mappings", "Complete function conversion across all source databases", GREEN),
    ("5", "OLTP→OLAP Patterns", "Rewrite procedural code to set-based SQL for columnar performance", TEAL),
    ("4", "UDx Languages", "C++, Python, Java, R for custom function development", PURPLE),
    ("4", "ML Categories", "Regression, classification, clustering, time series", RED),
    ("18", "Reference Guides", "Comprehensive documentation for every aspect", ACCENT_BLUE),
]

for i, (num, title, desc, color) in enumerate(takeaways):
    col = i % 3
    row = i // 3
    x = Inches(0.5 + col * 4.15)
    y = Inches(1.3 + row * 2.2)
    card = add_shape(slide, x, y, Inches(4.0), Inches(2.0), RGBColor(0x22, 0x33, 0x55), border_color=color, border_width=Pt(2))
    add_text_box(slide, x, y + Inches(0.2), Inches(4.0), Inches(0.6), num,
                font_size=36, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, y + Inches(0.85), Inches(4.0), Inches(0.4), title,
                font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), y + Inches(1.25), Inches(3.6), Inches(0.7),
                desc, font_size=11, color=RGBColor(0xAA, 0xBB, 0xCC), alignment=PP_ALIGN.CENTER)

add_shape(slide, Inches(0), Inches(6.8), Inches(13.333), Inches(0.7), MID_BLUE)
add_text_box(slide, Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.6),
             "The Vertica Expert Skill — Your Complete Guide to Database Migration, Development & Machine Learning",
             font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)


# ── Save ────────────────────────────────────────────────────────────────────

output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "Vertica_Expert_Skill_Overview.pptx")
prs.save(output_path)
print(f"✅ Presentation saved to: {output_path}")
print(f"   Total slides: {len(prs.slides)}")
